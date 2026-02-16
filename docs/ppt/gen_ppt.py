"""
LinkingChat PPT Generator — "Hermès Tech" Design System
Generates two presentations:
  A) Presentation_Overview.pptx  — Business / investor audience
  B) Presentation_Visuals.pptx   — CTO / tech review audience

Usage:
  python gen_ppt.py
"""

from __future__ import annotations

import os
from dataclasses import dataclass
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn


# ─────────────────────────────────────────────────────────────────
# 1. Design System: "Hermès Tech"
# ─────────────────────────────────────────────────────────────────


@dataclass(frozen=True)
class DesignSystem:
    """Strict color + typography tokens."""

    # ── Narrative Palette ──
    BG_VOID: RGBColor = RGBColor(0x0C, 0x13, 0x21)       # 极深藏青
    ORANGE: RGBColor = RGBColor(0xF3, 0x70, 0x21)         # Hermès Orange — The Solution
    RED: RGBColor = RGBColor(0xCC, 0x14, 0x2E)            # Rouge H — The Pain
    GREY: RGBColor = RGBColor(0x4A, 0x5A, 0x75)           # Titanium Grey — Structure
    TEXT: RGBColor = RGBColor(0xF0, 0xED, 0xE8)           # 暖白正文
    NOTE: RGBColor = RGBColor(0x8A, 0x9A, 0xB5)           # 灰蓝注释
    CARD_FILL: RGBColor = RGBColor(0x11, 0x1B, 0x2B)      # Ghost card 5% fill
    DARK_ACCENT: RGBColor = RGBColor(0x16, 0x22, 0x38)    # Subtle accent bg

    # ── Typography ──
    FONT_CN: str = "Microsoft YaHei"
    FONT_EN: str = "Arial"

    # ── Slide dimensions (16:9) ──
    WIDTH: int = Inches(13.333)
    HEIGHT: int = Inches(7.5)


DS = DesignSystem()


# ─────────────────────────────────────────────────────────────────
# 2. Helper utilities
# ─────────────────────────────────────────────────────────────────


def _set_slide_master_bg(prs: Presentation) -> None:
    """Lock background on every slide layout via the slide master."""
    master = prs.slide_masters[0]
    bg = master.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DS.BG_VOID

    for layout in prs.slide_layouts:
        bg2 = layout.background
        fill2 = bg2.fill
        fill2.solid()
        fill2.fore_color.rgb = DS.BG_VOID


def _set_slide_bg(slide) -> None:
    """Explicitly set individual slide background (belt-and-suspenders)."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DS.BG_VOID


def _add_slide(prs: Presentation):
    """Add a blank slide with background applied."""
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide)
    return slide


def _txbox(slide, left, top, width, height) -> "Shape":
    """Add a textbox at given position (all in Inches)."""
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )


def _set_font(run, size: int, color: RGBColor = DS.TEXT,
              bold: bool = False, italic: bool = False,
              font_name: Optional[str] = None):
    """Configure run font properties."""
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name or DS.FONT_EN
    # CJK east-asian font
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn("a:ea"), DS.FONT_CN if font_name is None else font_name)


def _add_text(slide, left, top, width, height, text: str,
              size: int = 14, color: RGBColor = DS.TEXT,
              bold: bool = False, align=PP_ALIGN.LEFT,
              font_name: Optional[str] = None):
    """Shortcut: add a single-paragraph textbox."""
    txBox = _txbox(slide, left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size, color, bold, font_name=font_name)
    return txBox


def _add_multiline(slide, left, top, width, height, lines: list[tuple],
                   line_spacing: float = 1.3, align=PP_ALIGN.LEFT):
    """Add textbox with multiple styled lines.
    lines: [(text, size, color, bold), ...]
    """
    txBox = _txbox(slide, left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(size * (line_spacing - 1) + 2)
        run = p.add_run()
        run.text = text
        _set_font(run, size, color, bold)
    return txBox


def _draw_line(slide, x1, y1, x2, y2,
               color: RGBColor = DS.GREY,
               width_pt: float = 1.0,
               dash: Optional[MSO_LINE_DASH_STYLE] = None):
    """Draw a connector line."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2),
    )
    ln = connector.line
    ln.color.rgb = color
    ln.width = Pt(width_pt)
    if dash:
        ln.dash_style = dash
    return connector


def _draw_rect(slide, left, top, width, height,
               fill_color: Optional[RGBColor] = None,
               border_color: RGBColor = DS.GREY,
               border_width: float = 1.0,
               border_dash: Optional[MSO_LINE_DASH_STYLE] = None,
               corner_radius: Optional[float] = None):
    """Draw a rectangle shape."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_width)
    if border_dash:
        shape.line.dash_style = border_dash
    return shape


def _draw_circle(slide, cx, cy, r,
                 fill_color: Optional[RGBColor] = None,
                 border_color: RGBColor = DS.ORANGE,
                 border_dash=MSO_LINE_DASH_STYLE.DASH,
                 border_width: float = 1.0):
    """Draw a circle (geometric stencil icon)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r),
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_width)
    if border_dash:
        shape.line.dash_style = border_dash
    return shape


def _corner_marks(slide, left, top, width, height,
                  size: float = 0.15,
                  color: RGBColor = DS.GREY,
                  width_pt: float = 0.75):
    """Draw L-shaped corner marks (camera viewfinder / blueprint aesthetic)."""
    r = left + width
    b = top + height
    s = size
    lines = [
        # Top-left
        (left, top, left + s, top), (left, top, left, top + s),
        # Top-right
        (r - s, top, r, top), (r, top, r, top + s),
        # Bottom-left
        (left, b - s, left, b), (left, b, left + s, b),
        # Bottom-right
        (r - s, b, r, b), (r, b - s, r, b),
    ]
    for x1, y1, x2, y2 in lines:
        _draw_line(slide, x1, y1, x2, y2, color=color, width_pt=width_pt)


def _ghost_card(slide, left, top, width, height,
                with_corners: bool = True):
    """Draw a ghost card (semi-transparent fill + border + corner marks)."""
    _draw_rect(slide, left, top, width, height,
               fill_color=DS.CARD_FILL,
               border_color=DS.GREY, border_width=0.75)
    if with_corners:
        _corner_marks(slide, left, top, width, height)


def _section_number(slide, num: int, top: float = 0.5):
    """Draw a section number like '01 /' in the top-left."""
    _add_text(slide, 0.6, top, 1.2, 0.4,
              f"{num:02d} /", size=11, color=DS.GREY, bold=False)


def _page_accent_line(slide):
    """Thin horizontal accent line near top."""
    _draw_line(slide, 0.6, 1.0, 3.6, 1.0,
               color=DS.ORANGE, width_pt=1.5)


def _arrow_right(slide, x1, y, x2,
                 color: RGBColor = DS.GREY,
                 dash=MSO_LINE_DASH_STYLE.DASH):
    """Draw a horizontal arrow (line + small triangle head)."""
    _draw_line(slide, x1, y, x2, y, color=color, width_pt=1.0, dash=dash)
    # arrowhead triangle
    sz = 0.08
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(x2 - 0.01), Inches(y - sz / 2), Inches(sz * 1.5), Inches(sz),
    )
    shape.rotation = 90
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _card_with_label(slide, left, top, w, h, label, sublabel="",
                     accent_color=DS.ORANGE, icon_type="circle"):
    """Ghost card with a geometric icon and label inside."""
    _ghost_card(slide, left, top, w, h)
    # Geometric icon
    cx = left + w / 2
    if icon_type == "circle":
        _draw_circle(slide, cx, top + 0.45, 0.22,
                     border_color=accent_color, border_dash=MSO_LINE_DASH_STYLE.DASH)
    elif icon_type == "rect":
        _draw_rect(slide, cx - 0.22, top + 0.23, 0.44, 0.44,
                   border_color=accent_color, border_width=1.0,
                   border_dash=MSO_LINE_DASH_STYLE.DASH)
    elif icon_type == "stack":
        for offset in [0, 0.08, 0.16]:
            _draw_rect(slide, cx - 0.2 + offset * 0.3, top + 0.25 + offset, 0.4, 0.28,
                       border_color=accent_color, border_width=0.75,
                       border_dash=MSO_LINE_DASH_STYLE.DASH)
    # Label
    _add_text(slide, left + 0.1, top + h * 0.5, w - 0.2, 0.35,
              label, size=12, color=DS.TEXT, bold=True, align=PP_ALIGN.CENTER)
    if sublabel:
        _add_text(slide, left + 0.1, top + h * 0.5 + 0.35, w - 0.2, 0.6,
                  sublabel, size=9, color=DS.NOTE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────
# 3. PPT A — Presentation_Overview  (Business / Investor)
# ─────────────────────────────────────────────────────────────────


def build_overview() -> None:
    prs = Presentation()
    prs.slide_width = DS.WIDTH
    prs.slide_height = DS.HEIGHT
    _set_slide_master_bg(prs)

    # ── Slide 1: Cover ──────────────────────────────────────────
    s = _add_slide(prs)
    # Thin decorative lines
    _draw_line(s, 0.6, 5.8, 4.0, 5.8, DS.ORANGE, 2.0)
    _draw_line(s, 0.6, 5.9, 2.5, 5.9, DS.GREY, 0.75, MSO_LINE_DASH_STYLE.DASH)
    # Title cluster — bottom-left
    _add_text(s, 0.6, 3.8, 8, 1.0,
              "LinkingChat", size=48, color=DS.TEXT, bold=True)
    _add_text(s, 0.6, 4.8, 10, 0.6,
              "AI-Native Social × Remote Control — 产品概览",
              size=20, color=DS.NOTE)
    _add_text(s, 0.6, 6.2, 6, 0.4,
              "Ghost Mate  |  2026", size=12, color=DS.GREY)
    # Corner marks on whole slide
    _corner_marks(s, 0.3, 0.3, 12.7, 6.9, size=0.25, color=DS.GREY, width_pt=0.5)

    # ── Slide 2: Context ────────────────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 1)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 8, 0.6,
              "行业现状  Context", size=28, color=DS.TEXT, bold=True)

    _add_multiline(s, 0.6, 2.2, 5.5, 4.0, [
        ("远程办公已成新常态", 16, DS.TEXT, True),
        ("但工具链严重碎片化：聊天、远程桌面、任务管理、AI 助手各自为政。", 13, DS.NOTE, False),
        ("", 8, DS.BG_VOID, False),
        ("效率瓶颈", 16, DS.TEXT, True),
        ("73% 的远程团队每天在 4+ 个工具之间切换，\n平均每次上下文切换损失 23 分钟专注时间。", 13, DS.NOTE, False),
        ("", 8, DS.BG_VOID, False),
        ("AI 浪潮下的机会", 16, DS.TEXT, True),
        ("大模型能力爆发，但缺乏真正的「执行层」——\nAI 能聊天，却不能帮你干活。", 13, DS.NOTE, False),
    ])

    # Right-side decorative card
    _ghost_card(s, 7.2, 2.0, 5.2, 4.5)
    _add_multiline(s, 7.5, 2.4, 4.6, 3.5, [
        ("现状痛点速览", 14, DS.GREY, True),
        ("", 6, DS.BG_VOID, False),
        ("◇  社交 + 办公工具割裂", 12, DS.NOTE, False),
        ("◇  AI 只能「建议」不能「执行」", 12, DS.NOTE, False),
        ("◇  远程控制缺乏安全审批机制", 12, DS.NOTE, False),
        ("◇  跨设备协作体验差", 12, DS.NOTE, False),
        ("◇  团队沟通与任务执行脱节", 12, DS.NOTE, False),
    ])

    # ── Slide 3: The Pain ───────────────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 2)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 8, 0.6,
              "核心痛点  The Pain", size=28, color=DS.RED, bold=True)

    pains = [
        ("工具碎片化", "聊天用微信 / Slack，远程用 TeamViewer，\nAI 用 ChatGPT——来回切换，上下文断裂"),
        ("AI 只说不做", "现有 AI 助手只能给建议，\n无法真正帮用户执行任务"),
        ("远程操作无安全审批", "直接远程控制=完全信任，\n缺乏「草稿->确认->执行」机制"),
    ]
    for i, (title, desc) in enumerate(pains):
        x = 0.6 + i * 4.1
        _ghost_card(s, x, 2.4, 3.7, 3.8)
        # Red accent dot
        _draw_circle(s, x + 0.35, 2.85, 0.12,
                     border_color=DS.RED, border_dash=MSO_LINE_DASH_STYLE.DASH)
        _add_text(s, x + 0.6, 2.65, 2.8, 0.4,
                  title, size=16, color=DS.RED, bold=True)
        _add_text(s, x + 0.3, 3.3, 3.1, 2.5,
                  desc, size=12, color=DS.NOTE)

    # ── Slide 4: The Concept ────────────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 3)
    _draw_line(s, 0.6, 3.0, 12.7, 3.0, DS.GREY, 0.5, MSO_LINE_DASH_STYLE.ROUND_DOT)
    _add_text(s, 0.6, 3.3, 12, 1.2,
              "聊天即指挥中心\n社交即生产力",
              size=40, color=DS.ORANGE, bold=True)
    _add_text(s, 0.6, 5.0, 10, 0.8,
              "LinkingChat = 即时通讯 × AI Agent × 远程执行\n"
              "一个对话窗口，完成从沟通到执行的全链路闭环。",
              size=16, color=DS.NOTE)
    _draw_line(s, 0.6, 6.0, 12.7, 6.0, DS.GREY, 0.5, MSO_LINE_DASH_STYLE.ROUND_DOT)

    # ── Slide 5: The Solution (3 pillars) ───────────────────────
    s = _add_slide(prs)
    _section_number(s, 4)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "核心能力  The Solution", size=28, color=DS.ORANGE, bold=True)

    pillars = [
        ("Draft & Verify", "代理草稿",
         "用户说意图 → AI 生成草稿\n→ 用户确认 → 才执行\nBot 永远不自主行动",
         "circle"),
        ("The Whisper", "耳语建议",
         "用户 @ai 触发 → 云端生成\n1 条最佳回复（预填输入框）\n+ 2 条备选方案",
         "rect"),
        ("Predictive Actions", "预测执行",
         "Bot 分析上下文 → 生成操作卡片\n危险命令自动拦截\n如：检测到编译错误 → 建议修复命令",
         "stack"),
    ]
    for i, (en, cn, desc, icon) in enumerate(pillars):
        x = 0.6 + i * 4.1
        _ghost_card(s, x, 2.2, 3.7, 4.5)
        _card_with_label(s, x + 0.2, 2.5, 3.3, 1.2, en, cn,
                         accent_color=DS.ORANGE, icon_type=icon)
        _add_text(s, x + 0.3, 4.0, 3.1, 2.5,
                  desc, size=11, color=DS.NOTE)

    # ── Slide 6: User Value ─────────────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 5)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "用户价值  User Value", size=28, color=DS.TEXT, bold=True)

    # Before / After comparison
    _ghost_card(s, 0.6, 2.2, 5.6, 4.5)
    _add_text(s, 0.9, 2.4, 3, 0.4,
              "BEFORE", size=14, color=DS.RED, bold=True)
    before_items = [
        "◇  4+ 个工具来回切换",
        "◇  AI 给建议，自己手动执行",
        "◇  远程桌面=全权限，无审批",
        "◇  聊天和任务执行完全脱节",
        "◇  上下文频繁丢失",
    ]
    for j, item in enumerate(before_items):
        _add_text(s, 0.9, 3.0 + j * 0.55, 5, 0.45,
                  item, size=12, color=DS.NOTE)

    _ghost_card(s, 7.0, 2.2, 5.6, 4.5)
    _add_text(s, 7.3, 2.4, 3, 0.4,
              "AFTER", size=14, color=DS.ORANGE, bold=True)
    after_items = [
        "◆  一个对话窗口，聊天 + 执行",
        "◆  AI 生成草稿，确认即执行",
        "◆  Draft & Verify 安全审批",
        "◆  从沟通到交付，零切换",
        "◆  全链路上下文自动保持",
    ]
    for j, item in enumerate(after_items):
        _add_text(s, 7.3, 3.0 + j * 0.55, 5, 0.45,
                  item, size=12, color=DS.TEXT)

    # Arrow between
    _arrow_right(s, 6.3, 4.45, 6.85, color=DS.ORANGE, dash=MSO_LINE_DASH_STYLE.DASH)

    # ── Slide 7: Architecture (simplified) ──────────────────────
    s = _add_slide(prs)
    _section_number(s, 6)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "产品架构  Cloud Brain + Local Hands",
              size=28, color=DS.TEXT, bold=True)

    # Three-tier boxes
    tiers = [
        (0.6, "Flutter 移动端", "社交界面\n发送指令\n确认草稿", DS.ORANGE),
        (4.6, "Cloud Brain", "NestJS 云服务\nWebSocket 网关\n意图规划 / LLM 路由\nAgent 逻辑", DS.ORANGE),
        (8.8, "Electron 桌面端", "社交 UI（类 Discord）\nOpenClaw Worker\nShell / 文件 / 自动化", DS.ORANGE),
    ]
    for x, title, desc, accent in tiers:
        _ghost_card(s, x, 2.4, 3.6, 3.8)
        _draw_circle(s, x + 1.8, 2.9, 0.25,
                     border_color=accent, border_dash=MSO_LINE_DASH_STYLE.DASH)
        _add_text(s, x + 0.2, 3.3, 3.2, 0.35,
                  title, size=14, color=DS.TEXT, bold=True, align=PP_ALIGN.CENTER)
        _add_text(s, x + 0.2, 3.75, 3.2, 2.2,
                  desc, size=11, color=DS.NOTE, align=PP_ALIGN.CENTER)

    # Arrows between tiers
    _arrow_right(s, 4.3, 4.3, 4.5, color=DS.GREY, dash=MSO_LINE_DASH_STYLE.DASH)
    _arrow_right(s, 8.3, 4.3, 8.7, color=DS.GREY, dash=MSO_LINE_DASH_STYLE.DASH)
    _add_text(s, 4.0, 4.55, 1.0, 0.3, "WSS", size=9, color=DS.GREY,
              align=PP_ALIGN.CENTER)
    _add_text(s, 8.2, 4.55, 1.0, 0.3, "WSS", size=9, color=DS.GREY,
              align=PP_ALIGN.CENTER)

    # ── Slide 8: First Milestone ────────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 7)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "首个里程碑  First Milestone", size=28, color=DS.TEXT, bold=True)

    _add_text(s, 0.6, 2.2, 11, 0.8,
              "手机发送一个干活指令 → 电脑直接干活 → 将结果发回手机端",
              size=18, color=DS.ORANGE, bold=True)

    # Flow steps
    steps = [
        ("01", "Mobile 发送\n工作指令", DS.ORANGE),
        ("02", "Cloud Brain\n解析意图", DS.GREY),
        ("03", "Desktop\n执行任务", DS.GREY),
        ("04", "结果回传\nMobile 确认", DS.ORANGE),
    ]
    for i, (num, label, accent) in enumerate(steps):
        x = 0.8 + i * 3.1
        _ghost_card(s, x, 3.5, 2.6, 2.5)
        _add_text(s, x + 0.1, 3.65, 0.6, 0.35,
                  num, size=22, color=accent, bold=True)
        _add_text(s, x + 0.2, 4.2, 2.2, 1.5,
                  label, size=13, color=DS.TEXT)
        if i < len(steps) - 1:
            _arrow_right(s, x + 2.65, 4.75, x + 3.0,
                         color=DS.GREY, dash=MSO_LINE_DASH_STYLE.DASH)

    # ── Slide 9: Roadmap ────────────────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 8)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "路线图  Roadmap", size=28, color=DS.TEXT, bold=True)

    phases = [
        ("Phase 0", "脚手架", "Monorepo 搭建\nCI/CD\n开发环境", "2 weeks"),
        ("Phase 1", "最小 PoC", "手机→云→桌面\n全链路贯通", "3 weeks"),
        ("Phase 2", "社交 MVP", "好友 / 群组\n消息 / 已读\nBot 框架", "5 weeks"),
        ("Phase 3", "AI 集成", "Whisper\nDraft & Verify\nPredictive Actions", "4 weeks"),
        ("Phase 4", "Polish", "性能优化\n安全审计\n公测准备", "3 weeks"),
    ]
    # Timeline line
    _draw_line(s, 0.8, 4.0, 12.5, 4.0, DS.GREY, 1.0, MSO_LINE_DASH_STYLE.DASH)

    for i, (phase, title, desc, dur) in enumerate(phases):
        x = 0.6 + i * 2.5
        # Dot on timeline
        _draw_circle(s, x + 0.5, 4.0, 0.08,
                     fill_color=DS.ORANGE if i <= 1 else DS.GREY,
                     border_color=DS.ORANGE if i <= 1 else DS.GREY,
                     border_dash=None)
        _add_text(s, x, 2.4, 2.2, 0.3,
                  phase, size=11, color=DS.ORANGE if i <= 1 else DS.GREY, bold=True)
        _add_text(s, x, 2.75, 2.2, 0.3,
                  title, size=13, color=DS.TEXT, bold=True)
        _add_text(s, x, 3.15, 2.2, 0.6,
                  dur, size=10, color=DS.NOTE)
        _add_text(s, x, 4.3, 2.2, 2.0,
                  desc, size=10, color=DS.NOTE)

    # ── Slide 10: Multi-Bot Architecture ────────────────────────
    s = _add_slide(prs)
    _section_number(s, 9)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "多 Bot 架构  Multi-Bot Framework",
              size=28, color=DS.TEXT, bold=True)

    _add_text(s, 0.6, 2.1, 11, 0.5,
              "MVP 即搭建多 Bot 框架，注册自动创建 Supervisor Bot + Coding Bot",
              size=14, color=DS.NOTE)

    bots = [
        ("Supervisor Bot", "通知聚合器\n智能管家\n所有 Bot 事件汇总\n不可删除，始终置顶", DS.ORANGE),
        ("Coding Bot", "远程执行代理\n代码 / Shell / 文件\nOpenClaw 集成\n默认置顶，可配置", DS.ORANGE),
        ("v1.x Bot 扩展", "社交媒体 Bot\n数据分析 Bot\n按需增加类型\n[待补充: 具体类型]", DS.GREY),
        ("v2.0 自定义", "用户自建 Bot\n开放创建能力\n自定义 Agent 配置\n[待补充: 开放策略]", DS.GREY),
    ]
    for i, (name, desc, accent) in enumerate(bots):
        x = 0.6 + i * 3.15
        _ghost_card(s, x, 2.8, 2.85, 3.8)
        _draw_circle(s, x + 1.42, 3.25, 0.22,
                     border_color=accent, border_dash=MSO_LINE_DASH_STYLE.DASH)
        _add_text(s, x + 0.15, 3.65, 2.55, 0.35,
                  name, size=12, color=DS.TEXT, bold=True, align=PP_ALIGN.CENTER)
        _add_text(s, x + 0.15, 4.1, 2.55, 2.2,
                  desc, size=10, color=DS.NOTE, align=PP_ALIGN.CENTER)

    # ── Slide 11: End ───────────────────────────────────────────
    s = _add_slide(prs)
    _draw_line(s, 0.6, 4.5, 6.0, 4.5, DS.ORANGE, 2.0)
    _add_text(s, 0.6, 2.5, 12, 1.0,
              "Chat is the new Terminal.",
              size=42, color=DS.ORANGE, bold=True)
    _add_text(s, 0.6, 3.6, 10, 0.6,
              "聊天即指挥中心，社交即生产力。",
              size=20, color=DS.NOTE)
    _add_text(s, 0.6, 5.0, 6, 0.4,
              "LinkingChat  ·  Ghost Mate  ·  2026",
              size=12, color=DS.GREY)
    _corner_marks(s, 0.3, 0.3, 12.7, 6.9, size=0.25, color=DS.GREY, width_pt=0.5)

    # ── Save ────────────────────────────────────────────────────
    out = os.path.join(os.path.dirname(__file__), "Presentation_Overview.pptx")
    prs.save(out)
    print(f"[OK] Saved: {out}")


# ─────────────────────────────────────────────────────────────────
# 4. PPT B — Presentation_Visuals  (CTO / Tech Review)
# ─────────────────────────────────────────────────────────────────


def build_visuals() -> None:
    prs = Presentation()
    prs.slide_width = DS.WIDTH
    prs.slide_height = DS.HEIGHT
    _set_slide_master_bg(prs)

    # ── Slide 1: Cover ──────────────────────────────────────────
    s = _add_slide(prs)
    _draw_line(s, 0.6, 5.8, 4.0, 5.8, DS.ORANGE, 2.0)
    _draw_line(s, 0.6, 5.9, 2.5, 5.9, DS.GREY, 0.75, MSO_LINE_DASH_STYLE.DASH)
    _add_text(s, 0.6, 3.5, 10, 1.0,
              "LinkingChat", size=48, color=DS.TEXT, bold=True)
    _add_text(s, 0.6, 4.5, 10, 0.5,
              "System Blueprints — Technical Review",
              size=22, color=DS.NOTE)
    _add_text(s, 0.6, 5.2, 6, 0.4,
              "Architecture · Data Flow · Protocol · UI Concept",
              size=12, color=DS.GREY)
    _add_text(s, 0.6, 6.2, 6, 0.4,
              "Ghost Mate  |  2026  |  CONFIDENTIAL",
              size=12, color=DS.GREY)
    _corner_marks(s, 0.3, 0.3, 12.7, 6.9, size=0.25, color=DS.GREY, width_pt=0.5)

    # ── Slide 2: Tech Stack Matrix ──────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 1)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 8, 0.6,
              "技术栈总览  Tech Stack", size=28, color=DS.TEXT, bold=True)

    stack_rows = [
        ("Layer",         "Technology",             "Rationale"),
        ("Cloud",         "NestJS / TypeScript",    "模块化、团队统一语言"),
        ("Database",      "PostgreSQL + Prisma",    "强类型 ORM、迁移管理"),
        ("Cache / PubSub","Redis",                  "Presence、消息广播"),
        ("Mobile",        "Flutter",                "iOS + Android 单代码库"),
        ("Desktop",       "Electron + Node.js",     "社交 UI + OpenClaw Worker"),
        ("Protocol",      "Socket.IO (WSS)",        "房间模型、自动重连"),
        ("Auth",          "JWT RS256",              "非对称签名、服务端验证"),
        ("LLM Routing",   "DeepSeek / Kimi 2.5",   "低成本模型 + 高能力模型路由"),
        ("Object Storage","MinIO (S3-compatible)",  "文件上传、头像、媒体"),
        ("Monorepo",      "Turborepo + pnpm",       "共享类型、统一构建"),
    ]
    # Table as ghost-card grid
    col_widths = [2.2, 3.2, 5.8]
    row_h = 0.38
    x_start = 0.6
    y_start = 2.2
    for ri, row in enumerate(stack_rows):
        y = y_start + ri * row_h
        for ci, cell in enumerate(row):
            x = x_start + sum(col_widths[:ci])
            w = col_widths[ci]
            if ri == 0:
                # Header row
                _draw_rect(s, x, y, w, row_h,
                           fill_color=DS.DARK_ACCENT,
                           border_color=DS.GREY, border_width=0.5)
                _add_text(s, x + 0.1, y + 0.02, w - 0.2, row_h - 0.04,
                          cell, size=10, color=DS.ORANGE, bold=True)
            else:
                _draw_rect(s, x, y, w, row_h,
                           border_color=DS.GREY, border_width=0.25)
                clr = DS.TEXT if ci == 1 else DS.NOTE
                _add_text(s, x + 0.1, y + 0.02, w - 0.2, row_h - 0.04,
                          cell, size=10, color=clr)

    # ── Slide 3: Architecture Layers ────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 2)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "系统架构  Architecture Layers", size=28, color=DS.TEXT, bold=True)

    layers = [
        ("Client Layer", 2.2, [
            ("Flutter Mobile", 0.6, 2.6, DS.ORANGE),
            ("Electron Desktop", 4.4, 2.6, DS.ORANGE),
            ("Web (future)", 8.2, 2.6, DS.GREY),
        ]),
        ("Gateway Layer", 3.6, [
            ("WebSocket Gateway", 0.6, 1.2, DS.ORANGE),
            ("REST API", 4.4, 1.2, DS.ORANGE),
            ("Auth Guard (JWT)", 8.2, 1.2, DS.GREY),
        ]),
        ("Service Layer", 5.0, [
            ("Chat Service", 0.6, 1.2, DS.ORANGE),
            ("User Service", 2.2, 1.2, DS.GREY),
            ("Bot Service", 3.8, 1.2, DS.GREY),
            ("AI / LLM Router", 5.4, 1.2, DS.ORANGE),
            ("OpenClaw Bridge", 7.0, 1.2, DS.GREY),
            ("File Service", 8.6, 1.2, DS.GREY),
        ]),
        ("Data Layer", 6.2, [
            ("PostgreSQL", 0.6, 2.6, DS.GREY),
            ("Redis", 4.4, 2.6, DS.GREY),
            ("MinIO (S3)", 8.2, 2.6, DS.GREY),
        ]),
    ]

    for layer_name, y, items in layers:
        # Layer label
        _add_text(s, 10.5, y, 2.5, 0.3,
                  layer_name, size=10, color=DS.NOTE, bold=True)
        _draw_line(s, 0.4, y + 0.35, 10.3, y + 0.35,
                   DS.GREY, 0.5, MSO_LINE_DASH_STYLE.ROUND_DOT)
        for name, x, w, accent in items:
            _draw_rect(s, x, y - 0.1, w, 0.35,
                       fill_color=DS.CARD_FILL,
                       border_color=accent, border_width=0.75,
                       border_dash=MSO_LINE_DASH_STYLE.DASH if accent == DS.GREY else None)
            _add_text(s, x + 0.05, y - 0.08, w - 0.1, 0.3,
                      name, size=9, color=DS.TEXT)

    # ── Slide 4: User Journey Flow ──────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 3)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "用户旅程  User Journey", size=28, color=DS.TEXT, bold=True)

    journey_steps = [
        ("用户发送\n意图消息", DS.ORANGE),
        ("Cloud Brain\n解析意图", DS.GREY),
        ("AI 生成\n执行草稿", DS.GREY),
        ("用户确认\n或修改", DS.ORANGE),
        ("Desktop\n执行命令", DS.GREY),
        ("结果回传\n+ 通知", DS.ORANGE),
    ]
    y_center = 3.8
    for i, (label, accent) in enumerate(journey_steps):
        x = 0.4 + i * 2.1
        # Rounded rect node
        _draw_rect(s, x, y_center - 0.5, 1.8, 1.2,
                   fill_color=DS.CARD_FILL,
                   border_color=accent, border_width=1.0,
                   corner_radius=0.15)
        _corner_marks(s, x, y_center - 0.5, 1.8, 1.2,
                      size=0.1, color=accent, width_pt=0.5)
        _add_text(s, x + 0.1, y_center - 0.3, 1.6, 0.8,
                  label, size=11, color=DS.TEXT, align=PP_ALIGN.CENTER)
        _add_text(s, x + 0.6, y_center - 0.85, 0.6, 0.3,
                  f"{i + 1:02d}", size=9, color=accent, bold=True,
                  align=PP_ALIGN.CENTER)
        if i < len(journey_steps) - 1:
            _arrow_right(s, x + 1.85, y_center, x + 2.05,
                         color=DS.GREY, dash=MSO_LINE_DASH_STYLE.DASH)

    # Bottom note
    _add_text(s, 0.6, 5.5, 11, 0.6,
              "Draft & Verify 模式：Bot 永远不自主执行——用户确认是必经节点",
              size=12, color=DS.NOTE)
    _draw_line(s, 0.6, 5.4, 12.5, 5.4, DS.GREY, 0.5, MSO_LINE_DASH_STYLE.ROUND_DOT)

    # ── Slide 5: Data Flow ──────────────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 4)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "数据流  Data Flow", size=28, color=DS.TEXT, bold=True)

    # Source → Process → Sink layout
    flow_cols = [
        ("SOURCE", 0.6, [
            ("Mobile Client", "circle"),
            ("Desktop Client", "circle"),
        ]),
        ("PROCESS", 4.5, [
            ("WS Gateway", "rect"),
            ("Chat Service", "rect"),
            ("LLM Router", "rect"),
            ("OpenClaw Bridge", "rect"),
        ]),
        ("SINK", 9.0, [
            ("PostgreSQL", "stack"),
            ("Redis PubSub", "stack"),
            ("MinIO Storage", "stack"),
        ]),
    ]
    for col_label, x, items in flow_cols:
        _add_text(s, x, 2.1, 3.5, 0.3,
                  col_label, size=10, color=DS.GREY, bold=True)
        _draw_line(s, x, 2.45, x + 3.5, 2.45,
                   DS.GREY, 0.5, MSO_LINE_DASH_STYLE.ROUND_DOT)
        for j, (name, icon) in enumerate(items):
            y = 2.7 + j * 1.1
            _ghost_card(s, x, y, 3.2, 0.85)
            icon_x = x + 0.35
            if icon == "circle":
                _draw_circle(s, icon_x, y + 0.42, 0.15,
                             border_color=DS.ORANGE, border_dash=MSO_LINE_DASH_STYLE.DASH)
            elif icon == "rect":
                _draw_rect(s, icon_x - 0.15, y + 0.25, 0.3, 0.3,
                           border_color=DS.GREY, border_width=0.75,
                           border_dash=MSO_LINE_DASH_STYLE.DASH)
            elif icon == "stack":
                for k in range(3):
                    _draw_rect(s, icon_x - 0.15 + k * 0.06, y + 0.2 + k * 0.08,
                               0.28, 0.2,
                               border_color=DS.GREY, border_width=0.5,
                               border_dash=MSO_LINE_DASH_STYLE.DASH)
            _add_text(s, x + 0.7, y + 0.2, 2.3, 0.45,
                      name, size=11, color=DS.TEXT)

    # Flow arrows
    _arrow_right(s, 3.9, 3.5, 4.4, color=DS.ORANGE, dash=MSO_LINE_DASH_STYLE.DASH)
    _arrow_right(s, 8.1, 3.5, 8.9, color=DS.ORANGE, dash=MSO_LINE_DASH_STYLE.DASH)

    # ── Slide 6: WebSocket Protocol ─────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 5)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "WebSocket 协议  Protocol Design", size=28, color=DS.TEXT, bold=True)

    # Namespace cards
    ns_data = [
        ("/chat", "社交命名空间", [
            "chat:send_message",
            "chat:message_received",
            "chat:typing",
            "chat:read",
            "chat:message_updated",
            "chat:message_deleted",
        ], DS.ORANGE),
        ("/device", "设备命名空间", [
            "device:execute_command",
            "device:command_result",
            "device:status_update",
            "device:heartbeat",
        ], DS.GREY),
        ("AI Events", "AI 事件", [
            "ai:whisper_request",
            "ai:whisper_response",
            "ai:draft_created",
            "ai:draft_confirmed",
            "ai:prediction_card",
        ], DS.ORANGE),
    ]
    for i, (ns, label, events, accent) in enumerate(ns_data):
        x = 0.6 + i * 4.2
        _ghost_card(s, x, 2.2, 3.8, 4.8)
        _add_text(s, x + 0.2, 2.35, 3.4, 0.35,
                  ns, size=14, color=accent, bold=True)
        _add_text(s, x + 0.2, 2.7, 3.4, 0.3,
                  label, size=10, color=DS.NOTE)
        _draw_line(s, x + 0.2, 3.05, x + 3.6, 3.05,
                   DS.GREY, 0.5, MSO_LINE_DASH_STYLE.ROUND_DOT)
        for j, evt in enumerate(events):
            _add_text(s, x + 0.3, 3.15 + j * 0.45, 3.2, 0.4,
                      evt, size=10, color=DS.TEXT,
                      font_name="Consolas")

    # Bottom: Room strategy
    _add_text(s, 0.6, 7.15, 12, 0.3,
              "Room Strategy:  u-{userId}  |  g-{groupId}  |  d-{deviceId}",
              size=11, color=DS.NOTE, font_name="Consolas")

    # ── Slide 7: Database Schema ────────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 6)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "数据库设计  Database Schema", size=28, color=DS.TEXT, bold=True)

    _add_text(s, 0.6, 2.0, 10, 0.3,
              "PostgreSQL + Prisma ORM  |  18 Models  |  CUID Keys  |  Soft Delete",
              size=11, color=DS.NOTE)

    models = [
        ("Core", ["User", "Account", "Session", "UserSetting"]),
        ("Social", ["Converse", "ConverseMember", "Friend", "FriendRequest"]),
        ("Messaging", ["Message", "MessageReaction", "MessageRead"]),
        ("Bot / AI", ["Bot", "BotConfig", "DraftState"]),
        ("Device", ["Device", "CommandLog"]),
        ("System", ["File", "Notification"]),
    ]
    for i, (group, tables) in enumerate(models):
        col = i % 3
        row = i // 3
        x = 0.6 + col * 4.2
        y = 2.6 + row * 2.4
        _ghost_card(s, x, y, 3.8, 2.0)
        _add_text(s, x + 0.15, y + 0.1, 3.5, 0.3,
                  group, size=11, color=DS.ORANGE, bold=True)
        _draw_line(s, x + 0.15, y + 0.4, x + 3.65, y + 0.4,
                   DS.GREY, 0.5, MSO_LINE_DASH_STYLE.ROUND_DOT)
        for j, tbl in enumerate(tables):
            _add_text(s, x + 0.2, y + 0.5 + j * 0.33, 3.4, 0.3,
                      f"◇  {tbl}", size=10, color=DS.TEXT)

    # ── Slide 8: UI Concept — Mobile ────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 7)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "移动端概念  Mobile UI Concept", size=28, color=DS.TEXT, bold=True)

    _add_text(s, 0.6, 2.0, 5, 0.3,
              "WeChat / WhatsApp 风格  ·  Less is More",
              size=12, color=DS.NOTE)

    # Phone frame (9:16 ratio)
    phone_w = 2.4
    phone_h = phone_w * 16 / 9
    phone_x = 1.5
    phone_y = 2.6
    _draw_rect(s, phone_x, phone_y, phone_w, phone_h,
               border_color=DS.GREY, border_width=1.5,
               corner_radius=0.2)
    _corner_marks(s, phone_x, phone_y, phone_w, phone_h,
                  size=0.15, color=DS.ORANGE, width_pt=0.75)

    # Status bar
    _draw_line(s, phone_x + 0.15, phone_y + 0.35,
               phone_x + phone_w - 0.15, phone_y + 0.35,
               DS.GREY, 0.5)
    _add_text(s, phone_x + 0.15, phone_y + 0.08, phone_w - 0.3, 0.25,
              "LinkingChat", size=8, color=DS.TEXT, bold=True,
              align=PP_ALIGN.CENTER)

    # Chat list items
    chat_items = [
        ("Supervisor Bot", "3 条新通知", DS.ORANGE),
        ("Coding Bot", "任务已完成 ✓", DS.ORANGE),
        ("张三", "好的，明天见", DS.GREY),
        ("开发群", "@你: PR 已合并", DS.GREY),
    ]
    for j, (name, preview, accent) in enumerate(chat_items):
        iy = phone_y + 0.5 + j * 0.65
        _draw_line(s, phone_x + 0.15, iy + 0.6,
                   phone_x + phone_w - 0.15, iy + 0.6,
                   DS.GREY, 0.25)
        # Avatar circle
        _draw_circle(s, phone_x + 0.4, iy + 0.3, 0.13,
                     border_color=accent,
                     border_dash=MSO_LINE_DASH_STYLE.DASH,
                     border_width=0.75)
        _add_text(s, phone_x + 0.6, iy + 0.05, 1.5, 0.25,
                  name, size=7, color=DS.TEXT, bold=True)
        _add_text(s, phone_x + 0.6, iy + 0.3, 1.5, 0.25,
                  preview, size=6, color=DS.NOTE)

    # Bottom tab bar
    _draw_line(s, phone_x + 0.15, phone_y + phone_h - 0.4,
               phone_x + phone_w - 0.15, phone_y + phone_h - 0.4,
               DS.GREY, 0.5)
    tabs = ["消息", "通讯录", "Bot", "我"]
    for j, tab in enumerate(tabs):
        tx = phone_x + 0.15 + j * (phone_w - 0.3) / 4
        _add_text(s, tx, phone_y + phone_h - 0.35, 0.5, 0.25,
                  tab, size=6, color=DS.ORANGE if j == 0 else DS.NOTE,
                  align=PP_ALIGN.CENTER)

    # Right side: Design notes
    _ghost_card(s, 5.5, 2.6, 7.0, 4.5)
    _add_multiline(s, 5.8, 2.8, 6.4, 4.0, [
        ("UI 设计原则", 14, DS.TEXT, True),
        ("", 6, DS.BG_VOID, False),
        ("01 /  Bot = 固定置顶系统联系人", 11, DS.NOTE, False),
        ("        类似微信「文件传输助手」", 9, DS.GREY, False),
        ("", 4, DS.BG_VOID, False),
        ("02 /  Supervisor Bot 聚合所有通知", 11, DS.NOTE, False),
        ("        普通聊天流 + BOT_NOTIFICATION 卡片", 9, DS.GREY, False),
        ("", 4, DS.BG_VOID, False),
        ("03 /  @ai 触发 Whisper，非自动推送", 11, DS.NOTE, False),
        ("        预填输入框 + ··· 展开备选", 9, DS.GREY, False),
        ("", 4, DS.BG_VOID, False),
        ("04 /  群聊中 @specificBot 直接调用", 11, DS.NOTE, False),
        ("        @ai = Supervisor 兜底", 9, DS.GREY, False),
        ("", 4, DS.BG_VOID, False),
        ("05 /  Draft & Verify 审批卡片式交互", 11, DS.NOTE, False),
        ("        草稿预览 → 确认/修改/拒绝", 9, DS.GREY, False),
    ])

    # ── Slide 9: OpenClaw Integration ───────────────────────────
    s = _add_slide(prs)
    _section_number(s, 8)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "OpenClaw 集成  Remote Execution", size=28, color=DS.TEXT, bold=True)

    _add_text(s, 0.6, 2.0, 11, 0.4,
              "OpenClaw = 开源 AI Agent Gateway（TypeScript, MIT）",
              size=13, color=DS.NOTE)

    # Integration diagram
    nodes = [
        (0.8, 3.2, 2.5, 1.5, "Cloud Brain\nNestJS", DS.ORANGE),
        (5.2, 3.2, 2.5, 1.5, "OpenClaw Node\n独立进程", DS.ORANGE),
        (9.5, 2.8, 2.8, 0.8, "Shell Executor", DS.GREY),
        (9.5, 3.8, 2.8, 0.8, "File I/O", DS.GREY),
        (9.5, 4.8, 2.8, 0.8, "Browser Auto", DS.GREY),
    ]
    for x, y, w, h, label, accent in nodes:
        _ghost_card(s, x, y, w, h)
        _add_text(s, x + 0.1, y + 0.15, w - 0.2, h - 0.3,
                  label, size=11, color=DS.TEXT, align=PP_ALIGN.CENTER)
        # Top accent bar
        _draw_line(s, x, y, x + w, y, color=accent, width_pt=1.5)

    # Arrows
    _arrow_right(s, 3.4, 3.95, 5.1, color=DS.ORANGE, dash=MSO_LINE_DASH_STYLE.DASH)
    _add_text(s, 3.6, 3.55, 1.3, 0.3, "WSS", size=9, color=DS.GREY,
              align=PP_ALIGN.CENTER)
    _arrow_right(s, 7.8, 3.2, 9.4, color=DS.GREY, dash=MSO_LINE_DASH_STYLE.DASH)
    _arrow_right(s, 7.8, 4.2, 9.4, color=DS.GREY, dash=MSO_LINE_DASH_STYLE.DASH)
    _arrow_right(s, 7.8, 5.2, 9.4, color=DS.GREY, dash=MSO_LINE_DASH_STYLE.DASH)

    # Security note
    _ghost_card(s, 0.6, 5.5, 11.8, 1.3)
    _add_multiline(s, 0.9, 5.6, 11.2, 1.0, [
        ("安全模型", 12, DS.RED, True),
        ("◇  OpenClaw Node 以独立进程运行，非 NestJS 子进程  |  "
         "◇  沙箱 + 白名单命令  |  "
         "◇  Draft & Verify 强制用户确认  |  "
         "◇  危险命令自动拦截标记",
         9, DS.NOTE, False),
    ])

    # ── Slide 10: Auth & Security ───────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 9)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "认证与安全  Auth & Security", size=28, color=DS.TEXT, bold=True)

    # Lock icon (geometric)
    lock_x, lock_y = 1.5, 2.6
    _draw_rect(s, lock_x, lock_y + 0.3, 0.8, 0.6,
               fill_color=DS.CARD_FILL,
               border_color=DS.ORANGE, border_width=1.5)
    _draw_circle(s, lock_x + 0.4, lock_y + 0.2, 0.25,
                 border_color=DS.ORANGE, border_dash=None, border_width=1.5)

    sec_items = [
        ("JWT RS256 非对称签名", "服务端持有私钥签发，客户端公钥验证\n支持密钥轮换，Token 7天有效 + Refresh Token"),
        ("WebSocket 认证", "连接时 handshake 携带 JWT\nSocket.IO middleware 统一拦截验证"),
        ("OpenClaw 安全沙箱", "独立进程隔离 + 命令白名单\n危险操作强制 Draft & Verify 审批"),
        ("数据安全", "Soft delete 保留审计轨迹\nJSONB 灵活扩展 + 强类型 Prisma 校验"),
    ]
    for i, (title, desc) in enumerate(sec_items):
        x = 3.0
        y = 2.4 + i * 1.2
        _ghost_card(s, x, y, 9.5, 1.0)
        _add_text(s, x + 0.2, y + 0.1, 3.5, 0.3,
                  title, size=12, color=DS.ORANGE, bold=True)
        _add_text(s, x + 0.2, y + 0.4, 9.0, 0.5,
                  desc, size=10, color=DS.NOTE)

    # ── Slide 11: Performance Targets ───────────────────────────
    s = _add_slide(prs)
    _section_number(s, 10)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "性能目标  Performance Targets", size=28, color=DS.TEXT, bold=True)

    perf_metrics = [
        ("< 2s", "消息镜像延迟\nMessage Mirror Latency", DS.ORANGE),
        ("< 3s", "远程执行延迟\nRemote Action Execution", DS.ORANGE),
        ("< 2s", "@ai 回复生成\nWhisper Generation", DS.ORANGE),
    ]
    for i, (val, label, accent) in enumerate(perf_metrics):
        x = 0.6 + i * 4.2
        _ghost_card(s, x, 2.4, 3.8, 3.0)
        _add_text(s, x + 0.3, 2.8, 3.2, 0.8,
                  val, size=40, color=accent, bold=True, align=PP_ALIGN.CENTER)
        _add_text(s, x + 0.3, 3.7, 3.2, 1.0,
                  label, size=13, color=DS.TEXT, align=PP_ALIGN.CENTER)

    # Scaling strategy
    _ghost_card(s, 0.6, 5.8, 12.0, 1.2)
    _add_multiline(s, 0.9, 5.9, 11.4, 1.0, [
        ("扩展策略", 12, DS.TEXT, True),
        ("水平扩展 (Horizontal Scaling) + Redis PubSub + Nginx LB  |  "
         "Rust 仅在数据证明热点瓶颈后考虑  |  "
         "架构解决扩展性，非语言",
         10, DS.NOTE, False),
    ])

    # ── Slide 12: Monorepo Structure ────────────────────────────
    s = _add_slide(prs)
    _section_number(s, 11)
    _page_accent_line(s)
    _add_text(s, 0.6, 1.2, 10, 0.6,
              "项目结构  Monorepo Layout", size=28, color=DS.TEXT, bold=True)

    # Tree visualization
    tree_data = [
        ("linkchat/", 0, DS.ORANGE),
        ("  apps/", 1, DS.TEXT),
        ("    server/          NestJS 云服务", 2, DS.NOTE),
        ("    web/             Web 客户端 (React)", 2, DS.NOTE),
        ("    desktop/         Electron 桌面端", 2, DS.NOTE),
        ("    mobile/          Flutter 移动端", 2, DS.NOTE),
        ("  packages/", 1, DS.TEXT),
        ("    shared/          共享类型定义", 2, DS.NOTE),
        ("    ws-protocol/     WebSocket 协议类型", 2, DS.NOTE),
        ("    api-client/      API 客户端 SDK", 2, DS.NOTE),
        ("    ui/              共享 UI 组件", 2, DS.NOTE),
        ("  docker/            Docker Compose 配置", 1, DS.NOTE),
        ("  prisma/            数据库 Schema + 迁移", 1, DS.NOTE),
        ("  turbo.json         Turborepo 配置", 1, DS.NOTE),
    ]
    _ghost_card(s, 0.6, 2.2, 6.5, 5.0)
    for i, (line, indent, color) in enumerate(tree_data):
        _add_text(s, 0.8, 2.35 + i * 0.33, 6.0, 0.3,
                  line, size=9, color=color, font_name="Consolas")

    # Right side: key principles
    _ghost_card(s, 7.5, 2.2, 5.2, 5.0)
    _add_multiline(s, 7.8, 2.4, 4.6, 4.6, [
        ("关键原则", 14, DS.TEXT, True),
        ("", 6, DS.BG_VOID, False),
        ("01 /  TypeScript Everywhere", 11, DS.ORANGE, True),
        ("        云 / 桌面 / Web 统一语言", 9, DS.NOTE, False),
        ("", 4, DS.BG_VOID, False),
        ("02 /  共享类型安全", 11, DS.ORANGE, True),
        ("        ws-protocol 包保证前后端类型一致", 9, DS.NOTE, False),
        ("", 4, DS.BG_VOID, False),
        ("03 /  Turborepo 增量构建", 11, DS.ORANGE, True),
        ("        缓存 + 并行，CI 速度提升 40%+", 9, DS.NOTE, False),
        ("", 4, DS.BG_VOID, False),
        ("04 /  Flutter 独立移动端", 11, DS.ORANGE, True),
        ("        iOS + Android 共享 95%+ 代码", 9, DS.NOTE, False),
        ("", 4, DS.BG_VOID, False),
        ("05 /  先 Monorepo，后拆分", 11, DS.ORANGE, True),
        ("        MVP 阶段保持简单，规模化后按需拆分", 9, DS.NOTE, False),
    ])

    # ── Slide 13: End ───────────────────────────────────────────
    s = _add_slide(prs)
    _draw_line(s, 0.6, 4.5, 6.0, 4.5, DS.ORANGE, 2.0)
    _add_text(s, 0.6, 2.2, 12, 0.8,
              "Cloud Brain + Local Hands",
              size=38, color=DS.ORANGE, bold=True)
    _add_text(s, 0.6, 3.2, 10, 0.6,
              "TypeScript Everywhere · Socket.IO · Prisma · OpenClaw",
              size=16, color=DS.NOTE)
    _add_text(s, 0.6, 3.9, 10, 0.5,
              "Architecture designed for iteration speed, not premature scale.",
              size=14, color=DS.GREY)
    _add_text(s, 0.6, 5.0, 6, 0.4,
              "LinkingChat  ·  Ghost Mate  ·  Technical Review  ·  2026",
              size=12, color=DS.GREY)
    _corner_marks(s, 0.3, 0.3, 12.7, 6.9, size=0.25, color=DS.GREY, width_pt=0.5)

    # ── Save ────────────────────────────────────────────────────
    out = os.path.join(os.path.dirname(__file__), "Presentation_Visuals.pptx")
    prs.save(out)
    print(f"[OK] Saved: {out}")


# ─────────────────────────────────────────────────────────────────
# 5. Main
# ─────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("=" * 50)
    print("LinkingChat PPT Generator — Hermès Tech")
    print("=" * 50)
    build_overview()
    build_visuals()
    print("\nDone. Files saved in docs/ppt/")
